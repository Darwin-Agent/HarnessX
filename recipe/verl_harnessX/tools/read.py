# Copyright 2026 Darwin-Agent
# SPDX-License-Identifier: MIT
"""
Read tool — read files from the local filesystem.

Supports text files, PDFs, DOCX, XLSX/CSV, and PPTX.
Read tool for verl_harnessX tool registry.
"""

from __future__ import annotations

import logging
import os
import time

from .base import tool

logger = logging.getLogger(__name__)

_success_count: int = 0
_fail_count: int = 0

_SCHEMA = {
    "type": "object",
    "properties": {
        "file_path": {
            "type": "string",
            "description": "The absolute path to the file to read",
        },
        "limit": {
            "type": "integer",
            "description": "Max number of lines to read (default 2000)",
        },
        "offset": {
            "type": "integer",
            "description": "Line number to start reading from (0-indexed)",
        },
        "pages": {
            "type": "string",
            "description": (
                "Page range for PDF files (e.g. '1-5', '3', '10-20'). "
                "Only applicable to .pdf files. Maximum 20 pages per request."
            ),
        },
    },
    "required": ["file_path"],
}

_OFFICE_EXTENSIONS = {".docx", ".xlsx", ".xls", ".csv", ".pptx"}

_MAX_PDF_PAGES = 20


def _get_sandbox():
    try:
        from harnessx.sandbox.base import get_current_sandbox

        return get_current_sandbox()
    except ImportError:
        return None


def _continuation_notice(start: int, end: int, total: int) -> str:
    return f"\n[Showing lines {start + 1}–{end} of {total} total. Use offset={end} to read the next section.]"


@tool(
    name="Read",
    description=(
        "Read a file from the local filesystem. Returns file contents with line numbers. "
        "Supports text files, PDFs, DOCX, XLSX/CSV spreadsheets, and PPTX presentations."
    ),
    input_schema=_SCHEMA,
)
async def read_tool(file_path: str, limit: int = 2000, offset: int = 0, pages: str | None = None) -> str:
    global _success_count, _fail_count
    t0 = time.monotonic()
    limit = int(limit) if limit is not None else 2000
    offset = int(offset) if offset is not None else 0

    if file_path.startswith(("http://", "https://")):
        _fail_count += 1
        logger.warning(
            "Read FAILED url-as-path: %s [total: %d ok, %d fail]",
            file_path[:80],
            _success_count,
            _fail_count,
        )
        return (
            f"Error: '{file_path}' is a URL, not a local file path. "
            "Use the WebFetch tool to retrieve web content, or use Bash with curl/wget to download the file first."
        )

    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        result = _read_pdf(file_path, pages)
        elapsed = time.monotonic() - t0
        if result.startswith("Error"):
            _fail_count += 1
            logger.warning(
                "Read FAILED (%.1fs) pdf %s [total: %d ok, %d fail]",
                elapsed,
                file_path[:80],
                _success_count,
                _fail_count,
            )
        else:
            _success_count += 1
            logger.warning(
                "Read OK (%.1fs) pdf %s [total: %d ok, %d fail]",
                elapsed,
                file_path[:80],
                _success_count,
                _fail_count,
            )
        return result
    if ext in _OFFICE_EXTENSIONS:
        result = _read_office(file_path, ext)
        elapsed = time.monotonic() - t0
        if result.startswith("Error"):
            _fail_count += 1
            logger.warning(
                "Read FAILED (%.1fs) %s %s [total: %d ok, %d fail]",
                elapsed,
                ext,
                file_path[:80],
                _success_count,
                _fail_count,
            )
        else:
            _success_count += 1
            logger.warning(
                "Read OK (%.1fs) %s %s [total: %d ok, %d fail]",
                elapsed,
                ext,
                file_path[:80],
                _success_count,
                _fail_count,
            )
        return result

    sandbox = _get_sandbox()
    resolved = sandbox.resolve(file_path) if sandbox is not None else file_path
    try:
        if sandbox is not None:
            content = await sandbox.read_file(resolved)
            lines = content.splitlines(keepends=True)
        else:
            with open(resolved, "r", encoding="utf-8", errors="replace") as f:
                lines = f.readlines()
        start = offset
        end = offset + limit if limit > 0 else len(lines)
        selected = lines[start:end]
        result = "".join(f"{start + i + 1}\t{line}" for i, line in enumerate(selected))
        if limit > 0 and end < len(lines):
            result += _continuation_notice(start, end, len(lines))
        elapsed = time.monotonic() - t0
        _success_count += 1
        logger.warning(
            "Read OK (%.1fs) %d lines %s [total: %d ok, %d fail]",
            elapsed,
            len(selected),
            file_path[:80],
            _success_count,
            _fail_count,
        )
        return result
    except FileNotFoundError:
        _fail_count += 1
        logger.warning(
            "Read FAILED file not found: %s [total: %d ok, %d fail]",
            file_path[:80],
            _success_count,
            _fail_count,
        )
        return f"Error: File not found: {file_path}"
    except Exception as e:
        _fail_count += 1
        elapsed = time.monotonic() - t0
        logger.warning(
            "Read FAILED (%.1fs) %s: %s [total: %d ok, %d fail]",
            elapsed,
            file_path[:80],
            e,
            _success_count,
            _fail_count,
        )
        return f"Error: {e}"


def _parse_page_range(pages: str, total: int) -> list[int]:
    pages = pages.strip()
    if "-" in pages:
        parts = pages.split("-", 1)
        start = max(1, int(parts[0].strip()))
        end = min(total, int(parts[1].strip()))
    else:
        start = end = int(pages.strip())
    if end - start + 1 > _MAX_PDF_PAGES:
        end = start + _MAX_PDF_PAGES - 1
    return list(range(start - 1, end))


def _read_pdf(file_path: str, pages: str | None) -> str:
    try:
        import pdfplumber
    except ImportError:
        return "Error: PDF reading requires: pip install pdfplumber"

    try:
        with pdfplumber.open(file_path) as pdf:
            total = len(pdf.pages)
            if pages:
                page_indices = _parse_page_range(pages, total)
            else:
                if total > _MAX_PDF_PAGES:
                    return (
                        f"Error: PDF has {total} pages. "
                        f"Provide the pages parameter to read specific pages "
                        f"(e.g. pages='1-{_MAX_PDF_PAGES}'). Maximum {_MAX_PDF_PAGES} pages per request."
                    )
                page_indices = list(range(total))

            parts: list[str] = []
            for i in page_indices:
                if i >= total:
                    continue
                text = pdf.pages[i].extract_text() or ""
                parts.append(f"--- Page {i + 1} ---\n{text}")
            return "\n\n".join(parts)
    except FileNotFoundError:
        return f"Error: File not found: {file_path}"
    except Exception as e:
        return f"Error reading PDF: {e}"


def _read_office(file_path: str, ext: str) -> str:
    try:
        if ext == ".docx":
            return _read_docx(file_path)
        elif ext in (".xlsx", ".xls"):
            return _read_spreadsheet(file_path, ext)
        elif ext == ".csv":
            return _read_csv(file_path)
        elif ext == ".pptx":
            return _read_pptx(file_path)
        else:
            return f"Error: Unsupported file type: {ext}"
    except FileNotFoundError:
        return f"Error: File not found: {file_path}"
    except ImportError as e:
        return f"Error: Missing dependency — {e}. pip install python-docx openpyxl python-pptx"
    except Exception as e:
        return f"Error reading {ext} file: {e}"


def _read_docx(file_path: str) -> str:
    import docx

    doc = docx.Document(file_path)
    parts: list[str] = []
    for i, para in enumerate(doc.paragraphs):
        text = para.text.strip()
        if text:
            parts.append(f"{i + 1}\t{text}")

    for t_idx, table in enumerate(doc.tables):
        parts.append(f"\n--- Table {t_idx + 1} ---")
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            parts.append("\t".join(cells))

    return "\n".join(parts) if parts else "(empty document)"


def _read_spreadsheet(file_path: str, ext: str) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    parts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"--- Sheet: {sheet_name} ---")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            parts.append("\t".join(cells))
    wb.close()
    return "\n".join(parts) if parts else "(empty workbook)"


def _read_csv(file_path: str) -> str:
    import csv

    parts: list[str] = []
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        for i, row in enumerate(reader):
            parts.append(f"{i + 1}\t" + "\t".join(row))
            if i >= 5000:
                parts.append(f"\n[... truncated at {i + 1} rows]")
                break
    return "\n".join(parts) if parts else "(empty file)"


def _read_pptx(file_path: str) -> str:
    from pptx import Presentation

    prs = Presentation(file_path)
    parts: list[str] = []
    for i, slide in enumerate(prs.slides):
        parts.append(f"--- Slide {i + 1} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
    return "\n".join(parts) if parts else "(empty presentation)"
